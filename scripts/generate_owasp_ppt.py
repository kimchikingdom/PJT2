from pathlib import Path
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


OUTPUT = Path("docs/ppt/OWASP_3Tier_Public_Health_Portal.pptx")


def set_run_style(run, size=20, bold=False, color=(15, 23, 42)):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)


def style_title(shape, text):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    set_run_style(run, size=30, bold=True)


def style_subtitle(shape, text):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_run_style(run, size=16, color=(71, 85, 105))


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    style_title(slide.shapes.title, title)
    style_subtitle(slide.placeholders[1], subtitle)


def add_bullet_slide(prs, title, bullets, source=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    style_title(slide.shapes.title, title)

    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()

    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        run = p.add_run()
        run.text = bullet
        set_run_style(run, size=18, color=(30, 41, 59))

    if source:
        box = slide.shapes.add_textbox(Inches(0.6), Inches(6.7), Inches(12.0), Inches(0.4))
        ft = box.text_frame
        ft.clear()
        p = ft.paragraphs[0]
        run = p.add_run()
        run.text = f"근거: {source}"
        set_run_style(run, size=11, color=(100, 116, 139))


def add_table_slide(prs, title, headers, rows, source=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(11.5), Inches(0.6))
    style_title(title_box, title)

    table_shape = slide.shapes.add_table(
        rows=len(rows) + 1,
        cols=len(headers),
        left=Inches(0.5),
        top=Inches(1.2),
        width=Inches(12.3),
        height=Inches(5.1),
    )
    table = table_shape.table

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                set_run_style(r, size=13, bold=True)

    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    set_run_style(run, size=12, color=(30, 41, 59))

    if source:
        box = slide.shapes.add_textbox(Inches(0.6), Inches(6.75), Inches(12.0), Inches(0.3))
        ft = box.text_frame
        ft.clear()
        p = ft.paragraphs[0]
        run = p.add_run()
        run.text = f"근거: {source}"
        set_run_style(run, size=11, color=(100, 116, 139))


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    today = date.today().isoformat()

    add_title_slide(
        prs,
        "OWASP Top 10 기반 3-Tier 웹서비스 보안 분석",
        f"공공 의료 민원 포털 프로젝트 발표 · {today}",
    )

    add_bullet_slide(
        prs,
        "발표 목표",
        [
            "3-Tier 구조(Web/WAS/DB)에서 취약점이 어떻게 발생하는지 프로젝트 기반으로 설명",
            "OWASP Top 10:2025 전 항목(A01~A10) 커버리지와 재현 근거 제시",
            "취약점 나열이 아닌 개선 우선순위와 운영 대책까지 연결",
        ],
        source="docs/IMPLEMENTATION_MASTER_PLAN.md, docs/OWASP_SCENARIOS.md",
    )

    add_bullet_slide(
        prs,
        "프로젝트 아키텍처 (3-Tier)",
        [
            "Web Tier: Nginx Reverse Proxy",
            "WAS Tier: Flask + SQLAlchemy (인증/민원/게시판/마이데이터 로직)",
            "DB Tier: MariaDB (업무 데이터 + 감사로그 + 마이데이터 스냅샷)",
            "요청 흐름: Client -> Nginx -> Flask -> MariaDB",
        ],
        source="docker-compose.yml, web/nginx.conf",
    )

    add_bullet_slide(
        prs,
        "서비스 범위",
        [
            "사용자 기능: 회원가입/로그인, 게시판, 민원, 공지, 마이페이지",
            "관리자 기능: 사용자·게시물·공지·민원 관리, 로그 모니터링",
            "의료 마이데이터: 제공기관 Mock API + DB 시드 데이터 기반 조회",
            "교육용 취약점 실습: /vulnlab, /security/scenarios",
        ],
        source="README.md, docs/FEATURE_MATRIX.md",
    )

    add_table_slide(
        prs,
        "OWASP Top 10:2025 커버리지",
        ["코드", "항목", "프로젝트 내 점검 포인트"],
        [
            ["A01", "Broken Access Control", "관리자/일반 사용자 경계, IDOR 차단"],
            ["A02", "Security Misconfiguration", "기본 설정/오류 노출 점검"],
            ["A03", "Software Supply Chain", "의존성 버전/패키지 검증"],
            ["A04", "Cryptographic Failures", "민감정보 전송/저장 보호"],
            ["A05", "Injection", "입력 처리와 쿼리 안전성"],
            ["A06", "Insecure Design", "업무 규칙/Rate limiting/상태 전이"],
            ["A07", "Authentication Failures", "로그인/세션 정책"],
            ["A08", "Integrity Failures", "데이터/코드 무결성"],
            ["A09", "Logging & Alerting Failures", "감사로그/탐지/경보"],
            ["A10", "Exceptional Conditions", "Fail-Open, 예외 처리"],
        ],
        source="docs/OWASP_SCENARIOS.md",
    )

    add_bullet_slide(
        prs,
        "A01~A05 핵심 요약",
        [
            "A01: 접근권한 누락 시 타인 데이터 노출 가능",
            "A02: 기본 설정/디버그 노출은 내부 구조 유출로 직결",
            "A03: 취약 패키지/공급망 관리 실패는 대규모 사고로 확장",
            "A04: 암호화/전송 보호 미흡 시 의료·개인정보 직접 노출",
            "A05: 입력값 처리 취약 시 데이터 조회/변조 리스크 발생",
        ],
        source="docs/OWASP_SCENARIOS.md, docs/API_SPEC.md",
    )

    add_bullet_slide(
        prs,
        "A06~A10 핵심 요약",
        [
            "A06: 설계 단계의 보안 규칙 누락은 구현 품질과 무관하게 취약",
            "A07: 인증 정책 부재(잠금/MFA/세션관리)는 계정 탈취 리스크 증가",
            "A08: 신뢰경계 검증 실패는 코드/데이터 위변조로 연결",
            "A09: 로그가 없거나 품질이 낮으면 사고 탐지·대응이 불가",
            "A10: 예외 처리 실패(Fail-Open)는 보호 로직 전체 우회로 이어짐",
        ],
        source="docs/OWASP_SCENARIOS.md",
    )

    add_bullet_slide(
        prs,
        "상세 사례 1: A01 접근제어",
        [
            "증상: 사용자 권한 분리 실패 시 타인 리소스 접근 가능",
            "영향: 개인정보·민원 데이터 무단 열람",
            "확인 포인트: /admin, /complaints/{id} 접근 통제",
            "개선: 리소스 소유권 검증 + 역할 기반 접근 제어 일관 적용",
        ],
        source="was/app/routes.py, docs/TEST_CASES.md",
    )

    add_bullet_slide(
        prs,
        "상세 사례 2: A05 인젝션",
        [
            "증상: 사용자 입력이 검증 없이 질의/명령 경계로 유입",
            "영향: 데이터 조회 우회, 정보 노출, 무결성 훼손",
            "확인 포인트: 검색/필터 입력 경로와 쿼리 실행 방식",
            "개선: 파라미터 바인딩, 입력 검증, 위험 패턴 차단",
        ],
        source="docs/OWASP_SCENARIOS.md, docs/API_SPEC.md",
    )

    add_bullet_slide(
        prs,
        "상세 사례 3: A07 인증실패",
        [
            "증상: 로그인 정책/세션 정책 부재 또는 취약",
            "영향: 계정 탈취, 권한 오남용",
            "확인 포인트: 실패 로그 누적, 잠금 정책, 세션 식별자 품질",
            "개선: 실패횟수 제한, 계정 잠금, MFA, 강한 세션 토큰",
        ],
        source="was/app/routes.py, docs/TEST_CASES.md",
    )

    add_bullet_slide(
        prs,
        "상세 사례 4: A09 로깅·알림",
        [
            "증상: 보안 이벤트가 누락되거나 알림 체계가 없음",
            "영향: 침해 탐지 지연, 사고 대응 실패",
            "확인 포인트: /admin/logs 이벤트 필터와 감사 로그 품질",
            "개선: 표준 이벤트 스키마, 민감정보 마스킹, 임계치 알림",
        ],
        source="was/app/templates/admin/logs.html, was/app/routes.py",
    )

    add_bullet_slide(
        prs,
        "상세 사례 5: A10 예외 처리",
        [
            "증상: 예외 상황에서 접근통제가 풀리거나 내부 정보가 노출",
            "영향: 보호 우회, 공격 표면 확대",
            "확인 포인트: Fail-Open/Fail-Closed 동작, 4xx/5xx 처리 일관성",
            "개선: Fail-Secure 기본값, 공통 에러 핸들러, 내부 정보 비노출",
        ],
        source="docs/OWASP_SCENARIOS.md, was/app/templates/errors/500.html",
    )

    add_bullet_slide(
        prs,
        "마이데이터 연동 관점 보안 포인트",
        [
            "현재 구조: 제공기관 Mock API + consent/token 기반 조회",
            "검증 포인트: client 인증, token 만료/폐기, 사용자-데이터 바인딩",
            "로그 포인트: mydata_fetch, mydata_report_download, web_request",
            "개선: 토큰 scope 분리, 재시도/이상패턴 탐지, 감사추적 강화",
        ],
        source="was/app/routes.py, was/app/provider_service.py",
    )

    add_table_slide(
        prs,
        "개선 로드맵 (우선순위)",
        ["구간", "목표", "핵심 액션"],
        [
            ["단기 (1~2주)", "고위험 차단", "접근제어 복구, 인증 정책 강화, 민감정보 마스킹"],
            ["중기 (3~4주)", "운영탐지 강화", "로그 표준화, 경보 룰, 취약점 스캔 자동화"],
            ["장기 (5주+)", "보안 내재화", "보안 테스트 CI 편입, 설계 리뷰 정례화"],
        ],
        source="docs/SERVICE_COMPLETION_PLAN.md, docs/TEST_RESULTS.md",
    )

    add_bullet_slide(
        prs,
        "평가/시연 체크리스트",
        [
            "A01~A10 각 항목에 대해 재현 근거 1개 이상 확보",
            "관리자 로그에서 관련 이벤트 추적 가능 여부 확인",
            "개선안 적용 전/후 비교(리스크 감소) 명확히 제시",
            "문서-코드-시연 흐름 불일치 0건 확인",
        ],
        source="docs/TEST_CASES.md, docs/FEATURE_MATRIX.md",
    )

    add_bullet_slide(
        prs,
        "결론",
        [
            "이 프로젝트는 3-Tier 구조에서 보안 실패 지점을 체계적으로 학습하기에 적합",
            "OWASP Top 10 전 항목을 실제 기능과 연결해 분석 가능한 상태",
            "발표 핵심은 '취약점 개수' + '근거 기반 분석' + '실행 가능한 개선안'",
        ],
        source="docs/OWASP_SCENARIOS.md, docs/PPT_PRODUCTION_PLAN.md",
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    out = build_presentation()
    print(out)
